import io
import json
import os
from collections import Counter
from typing import Any, Dict, Iterable, List, Tuple
from xml.sax.saxutils import escape

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, PieChart, Reference
except ImportError:  # pragma: no cover
    Workbook = None
    BarChart = None
    PieChart = None
    Reference = None


REPORT_DIR = "storage/reports"
ASSET_DIR = os.path.join(REPORT_DIR, "assets")
os.makedirs(ASSET_DIR, exist_ok=True)

SYSTEM_OWNERS = {"", "unknown", "postgres", "system", "default"}
PRIORITY_ORDER = {"critical": 0, "improvement": 1, "observation": 2}
OWNER_LABEL = "Data Engineering / Data Steward"


def map_owner(owner: str) -> str:
    return OWNER_LABEL


def deduplicate_items(items: Iterable[Dict[str, Any]]) -> List[Dict[str, Any]]:
    deduped: List[Dict[str, Any]] = []
    seen = set()

    for item in items or []:
        if not isinstance(item, dict):
            continue
        summary = str(item.get("summary") or item.get("action") or item.get("description") or "").strip().lower()
        owner = str(item.get("owner") or "").strip().lower()
        category = str(item.get("category") or item.get("priority") or item.get("risk") or "").strip().lower()

        if "mask" in summary or "token" in summary or "restrict access" in summary:
            summary = "apply masking or tokenization to protect sensitive data"
        elif "sample" in summary or "profil" in summary:
            summary = "improve profiling coverage and resolve sampling limitations"

        key = (summary, owner, category)
        if key in seen:
            continue
        seen.add(key)

        merged = dict(item)
        if "mask" in summary:
            merged["action"] = "Apply masking or tokenization to protect sensitive data"
            merged["summary"] = merged["action"]
        elif "profil" in summary or "sample" in summary:
            merged["action"] = "Increase sample coverage and resolve profiling limitations"
            merged["summary"] = merged["action"]
        deduped.append(merged)

    return deduped


def _business_masking_label(value: str) -> str:
    if value == "NOT_MASKED":
        return "Unprotected sensitive data"
    if value == "PARTIALLY_MASKED":
        return "Partially protected sensitive data"
    if value == "MASKED":
        return "Protected sensitive data"
    return "Masking not applicable"


def _business_risk_label(value: str) -> str:
    lowered = str(value or "").strip().lower()
    if lowered == "critical":
        return "Critical compliance risk"
    if lowered == "high":
        return "High compliance risk"
    if lowered == "medium":
        return "Moderate compliance risk"
    return "Low compliance risk"


def _business_status_label(value: str) -> str:
    if value == "ACTION_REQUIRED":
        return "Immediate remediation required"
    return str(value or "Status unavailable")


def _build_pii_findings(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for table, columns in (result.get("column_intelligence") or {}).get("tables", {}).items():
        for column in columns:
            if not column.get("pii_detected"):
                continue
            findings.append({
                "table": table,
                "column": column.get("column"),
                "pii_type": column.get("pii_type"),
                "masking": _business_masking_label(column.get("masking_status")),
                "risk": _business_risk_label(column.get("risk")),
                "assigned_owner": map_owner(column.get("owner")),
                "recommendation": (
                    "Apply masking or tokenization to protect sensitive data and restrict unauthorized access"
                    if column.get("masking_status") == "NOT_MASKED"
                    else "Maintain existing controls and continue monitoring"
                ),
                "impact": (
                    "Potential DPDP non-compliance and exposure of personal data"
                    if column.get("masking_status") == "NOT_MASKED"
                    else "Reduced exposure, though continued monitoring is recommended"
                ),
            })
    return findings


def _build_risks(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    risks = []
    for item in (result.get("raid") or {}).get("risks", []):
        risks.append({
            "summary": item.get("description"),
            "risk": _business_risk_label(item.get("severity")),
            "assigned_owner": map_owner(item.get("owner")),
            "impact": "Potential DPDP non-compliance and exposure of personal data",
            "category": "critical" if str(item.get("severity")).lower() in {"critical", "high"} else "improvement",
        })
    return deduplicate_items(risks)


def _build_issues(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    issues = []
    for item in (result.get("gap_analysis") or {}).get("gaps", []):
        issues.append({
            "summary": item.get("description"),
            "category": item.get("type"),
            "priority": "critical" if str(item.get("severity")).lower() == "high" else "improvement",
            "assigned_owner": map_owner(item.get("owner")),
            "impact": (
                "Potential DPDP non-compliance and exposure of personal data"
                if str(item.get("severity")).lower() == "high"
                else "May reduce profiling confidence and governance clarity"
            ),
        })
    for item in (result.get("raid") or {}).get("issues", []):
        issues.append({
            "summary": item.get("description"),
            "category": item.get("type"),
            "priority": "critical" if str(item.get("severity")).lower() == "high" else "improvement",
            "assigned_owner": map_owner(item.get("owner")),
            "impact": (
                "Potential DPDP non-compliance and exposure of personal data"
                if str(item.get("severity")).lower() == "high"
                else "May reduce profiling confidence and governance clarity"
            ),
        })
    return deduplicate_items(issues)


def _build_recommendations(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    recommendations = []
    for item in result.get("recommendations", []):
        recommendations.append({
            "action": item.get("action"),
            "summary": item.get("action"),
            "priority": "critical" if str(item.get("priority")).lower() == "high" else "improvement",
            "assigned_owner": map_owner(item.get("owner")),
            "category": "critical" if str(item.get("priority")).lower() == "high" else "improvement",
        })

    if any("mask" in str(item.get("action", "")).lower() for item in recommendations):
        recommendations.append({
            "action": "Apply masking or tokenization to protect sensitive data and restrict unauthorized access",
            "summary": "Apply masking or tokenization to protect sensitive data and restrict unauthorized access",
            "priority": "critical",
            "assigned_owner": OWNER_LABEL,
            "category": "critical",
        })

    return deduplicate_items(recommendations)


def _build_document_alignment(result: Dict[str, Any]) -> List[str]:
    return [
        f"{item.get('doc_id')}: {item.get('field')} expected {item.get('expected')} but actual is {item.get('actual')} ({item.get('status')})"
        for item in result.get("document_alignment", [])[:10]
    ]


def _build_priorities(pii_findings: List[Dict[str, Any]], issues: List[Dict[str, Any]], relationships: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
    critical = []
    improvements = []
    observations = []

    for item in pii_findings:
        if item["masking"] == "Unprotected sensitive data":
            critical.append({
                "summary": f"{item['table']}.{item['column']} contains {item['pii_type']} and is unprotected sensitive data.",
                "impact": item["impact"],
                "assigned_owner": item["assigned_owner"],
            })

    for item in issues:
        bucket = critical if item["priority"] == "critical" else improvements
        bucket.append({
            "summary": item["summary"],
            "impact": item["impact"],
            "assigned_owner": item["assigned_owner"],
        })

    for relationship in relationships:
        observations.append({
            "summary": f"{relationship.get('from_table')} -> {relationship.get('to_table')} relationship signal detected.",
            "impact": "Useful for lineage validation and access review",
            "assigned_owner": OWNER_LABEL,
        })

    return {
        "critical": deduplicate_items(critical),
        "improvements": deduplicate_items(improvements),
        "observations": deduplicate_items(observations),
    }


def build_report_context(status_payload: Dict[str, Any], job_id: str) -> Dict[str, Any]:
    result = status_payload.get("result") or {}
    pii_findings = _build_pii_findings(result)
    risks = _build_risks(result)
    issues = _build_issues(result)
    recommendations = _build_recommendations(result)
    ai_insights = deduplicate_items([
        {
            "summary": item.get("insight"),
            "severity": item.get("severity"),
            "type": item.get("type"),
        }
        for item in result.get("ai_insights", [])
    ])
    relationships = result.get("relationship_inference", []) or []
    document_alignment = _build_document_alignment(result)
    priorities = _build_priorities(pii_findings, issues, relationships)

    executive = (
        "This scan identified unmasked personal data in customer tables, posing a high compliance risk under DPDP. "
        "Immediate remediation is required before further data usage."
        if any(item["masking"] == "Unprotected sensitive data" for item in pii_findings)
        else "This scan did not identify critical personal-data exposure, but governance controls should continue to be monitored."
    )

    return {
        "job_id": job_id,
        "run_id": result.get("run_id", job_id),
        "summary": {
            "database": result.get("source_inventory", {}).get("database", "Unknown"),
            "final_score": result.get("scores", {}).get("final_score"),
            "quality_score": result.get("quality_score", result.get("scores", {}).get("quality_score")),
            "status": _business_status_label(result.get("dpdp_compliance", {}).get("status")),
            "executive_narrative": executive,
        },
        "pii_findings": pii_findings,
        "risks": risks,
        "issues": issues,
        "recommendations": recommendations,
        "ai_insights": ai_insights,
        "priorities": priorities,
        "data_quality": {
            "summary": result.get("dq_report", {}).get("summary", "Profiling insights were generated from sampled source data."),
            "dimensions": result.get("dq_report", {}).get("quality_dimensions", {}),
            "issues": result.get("dq_report", {}).get("table_wise_issues", []),
        },
        "appendix": {
            "relationships": relationships,
            "document_alignment": document_alignment,
            "workflow": result.get("workflow", []),
            "tables": result.get("source_inventory", {}).get("tables", []),
        },
        "raw_status": status_payload,
    }


def _chart_path(job_id: str, name: str) -> str:
    return os.path.join(ASSET_DIR, f"{job_id}_{name}.png")


def generate_visuals(context: Dict[str, Any]) -> Dict[str, str]:
    risk_counts = Counter(item.get("risk", "Low compliance risk") for item in context["risks"])
    pii_counts = Counter(item.get("masking", "Masking not applicable") for item in context["pii_findings"])
    table_counts = Counter(item.get("table") for item in context["pii_findings"])

    dq_issues = context["data_quality"]["issues"]
    dq_counts = Counter(item.get("severity", "medium") for item in dq_issues)

    visuals = {}

    def save_pie(path: str, title: str, data: List[Tuple[str, int]], colors_list: List[str]) -> None:
        fig, ax = plt.subplots(figsize=(7, 4.5))
        positive_data = [(label, value) for label, value in data if value > 0]
        if positive_data:
            labels = [label for label, _ in positive_data]
            values = [value for _, value in positive_data]
        else:
            labels = ["No findings"]
            values = [1]
        ax.pie(values, labels=labels, autopct="%1.0f%%", colors=colors_list[:len(values)], startangle=90)
        ax.set_title(title)
        fig.tight_layout()
        fig.savefig(path, dpi=160)
        plt.close(fig)

    def save_bar(path: str, title: str, data: List[Tuple[str, int]], color: str) -> None:
        fig, ax = plt.subplots(figsize=(8, 4.5))
        labels = [label for label, _ in data] or ["No data"]
        values = [value for _, value in data] or [0]
        ax.bar(labels, values, color=color)
        ax.set_title(title)
        ax.tick_params(axis="x", rotation=20)
        fig.tight_layout()
        fig.savefig(path, dpi=160)
        plt.close(fig)

    risk_path = _chart_path(context["job_id"], "risk")
    save_pie(
        risk_path,
        "Risk Distribution",
        [
            ("Critical", risk_counts.get("Critical compliance risk", 0)),
            ("High", risk_counts.get("High compliance risk", 0)),
            ("Medium", risk_counts.get("Moderate compliance risk", 0)),
            ("Low", risk_counts.get("Low compliance risk", 0)),
        ],
        ["#7f1d1d", "#ba4538", "#d18a1c", "#2f7d5a"],
    )
    visuals["risk"] = risk_path

    pii_path = _chart_path(context["job_id"], "pii")
    save_pie(
        pii_path,
        "PII Exposure",
        [
            ("Unprotected", pii_counts.get("Unprotected sensitive data", 0)),
            ("Partial", pii_counts.get("Partially protected sensitive data", 0)),
            ("Protected", pii_counts.get("Protected sensitive data", 0)),
        ],
        ["#ba4538", "#d18a1c", "#2f7d5a"],
    )
    visuals["pii"] = pii_path

    dq_path = _chart_path(context["job_id"], "dq")
    save_bar(
        dq_path,
        "Data Quality Signals",
        [
            ("High", dq_counts.get("high", 0)),
            ("Medium", dq_counts.get("medium", 0)),
            ("Low", dq_counts.get("low", 0)),
        ],
        "#0f6c7a",
    )
    visuals["dq"] = dq_path

    coverage_path = _chart_path(context["job_id"], "coverage")
    save_bar(
        coverage_path,
        "Table Coverage",
        list(table_counts.items()) or [("No tables", 0)],
        "#7fb3bd",
    )
    visuals["coverage"] = coverage_path
    return visuals


def _section_heading(text: str, styles):
    return Paragraph(text, styles["Heading2"])


def _paragraph_cell(value: Any, style) -> Paragraph:
    return Paragraph(escape(str(value or "")), style)


def _section_lines(title: str, items: Iterable[str]) -> List[str]:
    lines = ["", f"## {title}"]
    values = [item for item in items if item]
    lines.extend(values or ["No findings were identified."])
    return lines


def export_raw_json(job_id: str, status_payload: Dict[str, Any]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.json")
    with open(path, "w", encoding="utf-8") as handle:
        json.dump(status_payload, handle, separators=(",", ":"))
    return path


def export_markdown(job_id: str, context: Dict[str, Any]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.md")
    lines = [
        "# Data Governance Report",
        "",
        "## Executive Summary",
        context["summary"]["executive_narrative"],
        "",
        "## Risk Overview",
    ]
    lines.extend(f"- {item['risk']}: {item['summary']}" for item in context["risks"])
    lines.extend(["", "## PII Exposure", "| Table | Column | Exposure | Risk |", "| --- | --- | --- | --- |"])
    lines.extend(
        f"| {item['table']} | {item['column']} | {item['masking']} | {item['risk']} |"
        for item in context["pii_findings"]
    )
    lines.extend(_section_lines("Data Quality", [context["data_quality"]["summary"]]))
    lines.extend(_section_lines("AI Insights", [f"- {item['summary']}" for item in context["ai_insights"]]))
    lines.extend(["", "## Recommendations"])
    lines.extend(f"- {item['action']}" for item in context["recommendations"])
    with open(path, "w", encoding="utf-8") as handle:
        handle.write("\n".join(lines))
    return path


def export_text(job_id: str, context: Dict[str, Any]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.txt")
    lines = [
        "DATA GOVERNANCE REPORT",
        "======================",
        "",
        "EXECUTIVE SUMMARY",
        context["summary"]["executive_narrative"],
        "",
        "RISK OVERVIEW",
    ]
    lines.extend(f"- {item['risk']}: {item['summary']}" for item in context["risks"])
    lines.extend([
        "",
        "PII EXPOSURE",
    ])
    lines.extend(f"- {item['table']}.{item['column']}: {item['masking']} ({item['risk']})" for item in context["pii_findings"])
    lines.extend([
        "",
        "DATA QUALITY",
        context["data_quality"]["summary"],
        "",
        "AI INSIGHTS",
    ])
    lines.extend(f"- {item['summary']}" for item in context["ai_insights"])
    lines.extend([
        "",
        "RECOMMENDATIONS",
    ])
    lines.extend(f"- {item['action']}" for item in context["recommendations"])
    with open(path, "w", encoding="utf-8") as handle:
        handle.write("\n".join(lines))
    return path


def export_docx(job_id: str, context: Dict[str, Any], visuals: Dict[str, str]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.docx")
    if Document is None:
        raise RuntimeError("python-docx is required for DOCX export")

    document = Document()
    document.add_heading("Data Governance Report", 0)
    document.add_heading("Executive Summary", level=1)
    document.add_paragraph(context["summary"]["executive_narrative"])

    document.add_heading("Risk Overview", level=1)
    document.add_picture(visuals["risk"], width=Inches(6.2))
    for item in context["risks"][:5]:
        document.add_paragraph(item["summary"], style="List Bullet")

    document.add_heading("PII Exposure", level=1)
    document.add_picture(visuals["pii"], width=Inches(6.2))
    for item in context["pii_findings"][:5]:
        document.add_paragraph(f"{item['table']}.{item['column']}: {item['masking']}", style="List Bullet")

    document.add_heading("Data Quality", level=1)
    document.add_picture(visuals["dq"], width=Inches(6.2))
    document.add_paragraph(context["data_quality"]["summary"])

    document.add_heading("AI Insights", level=1)
    for item in context["ai_insights"][:5]:
        document.add_paragraph(item["summary"], style="List Bullet")

    document.add_heading("Recommendations", level=1)
    for item in context["recommendations"][:5]:
        document.add_paragraph(item["action"], style="List Bullet")

    document.save(path)
    return path


def export_xlsx(job_id: str, context: Dict[str, Any], visuals: Dict[str, str]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.xlsx")
    if Workbook is None:
        raise RuntimeError("openpyxl is required for Excel export")

    workbook = Workbook()
    summary_ws = workbook.active
    summary_ws.title = "Summary"
    summary_ws.append(["Metric", "Value"])
    summary_ws.append(["Run ID", context["run_id"]])
    summary_ws.append(["Database", context["summary"]["database"]])
    summary_ws.append(["Final Score", context["summary"]["final_score"]])
    summary_ws.append(["Compliance Status", context["summary"]["status"]])

    pii_ws = workbook.create_sheet("PII Findings")
    pii_ws.append(["Table", "Column", "Exposure", "Risk", "Assigned Owner"])
    for item in context["pii_findings"]:
        pii_ws.append([item["table"], item["column"], item["masking"], item["risk"], item["assigned_owner"]])

    risk_ws = workbook.create_sheet("Risks")
    risk_ws.append(["Risk", "Assigned Owner", "Impact"])
    for item in context["risks"]:
        risk_ws.append([item["summary"], item["assigned_owner"], item["impact"]])

    rec_ws = workbook.create_sheet("Recommendations")
    rec_ws.append(["Action", "Priority", "Assigned Owner"])
    for item in context["recommendations"]:
        rec_ws.append([item["action"], item["priority"], item["assigned_owner"]])

    profiling_ws = workbook.create_sheet("Data Quality")
    profiling_ws.append(["Table", "Issue", "Severity"])
    for item in context["data_quality"]["issues"]:
        profiling_ws.append([item.get("table"), item.get("issue"), item.get("severity")])

    charts_ws = workbook.create_sheet("Charts")
    charts_ws.append(["Risk Level", "Count"])
    risk_rows = [
        ("Critical", sum(1 for item in context["risks"] if item["risk"] == "Critical compliance risk")),
        ("High", sum(1 for item in context["risks"] if item["risk"] == "High compliance risk")),
        ("Medium", sum(1 for item in context["risks"] if item["risk"] == "Moderate compliance risk")),
        ("Low", sum(1 for item in context["risks"] if item["risk"] == "Low compliance risk")),
    ]
    for row in risk_rows:
        charts_ws.append(list(row))
    charts_ws.append([])
    charts_ws.append(["DQ Severity", "Count"])
    dq_counts = Counter(item.get("severity", "medium") for item in context["data_quality"]["issues"])
    dq_start = charts_ws.max_row + 1
    for row in [
        ("High", dq_counts.get("high", 0)),
        ("Medium", dq_counts.get("medium", 0)),
        ("Low", dq_counts.get("low", 0)),
    ]:
        charts_ws.append(list(row))

    if PieChart and Reference:
        pie = PieChart()
        labels = Reference(charts_ws, min_col=1, min_row=2, max_row=5)
        data = Reference(charts_ws, min_col=2, min_row=1, max_row=5)
        pie.add_data(data, titles_from_data=True)
        pie.set_categories(labels)
        pie.title = "Risk Distribution"
        charts_ws.add_chart(pie, "E2")

        bar = BarChart()
        bar_data = Reference(charts_ws, min_col=2, min_row=dq_start - 1, max_row=dq_start + 2)
        bar_categories = Reference(charts_ws, min_col=1, min_row=dq_start, max_row=dq_start + 2)
        bar.add_data(bar_data, titles_from_data=True)
        bar.set_categories(bar_categories)
        bar.title = "Data Quality Issues"
        charts_ws.add_chart(bar, "E18")

    workbook.save(path)
    return path


def export_pptx(job_id: str, context: Dict[str, Any], visuals: Dict[str, str]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.pptx")
    presentation = Presentation()

    def add_bullets(slide, items: List[str], top: float = 1.4, left: float = 0.7, width: float = 8.2, height: float = 4.8) -> None:
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        bullet_items = [str(item) for item in items[:5] if item] or ["No findings were identified."]
        for index, item in enumerate(bullet_items):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Overview"
    add_bullets(slide, [
        context["summary"]["executive_narrative"],
        f"Final score: {context['summary']['final_score']}",
        f"Quality score: {context['summary']['quality_score']}",
        f"Status: {context['summary']['status']}",
    ])

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Risk Summary"
    slide.shapes.add_picture(visuals["risk"], Inches(0.6), Inches(1.2), width=Inches(5.6))
    add_bullets(slide, [item["summary"] for item in context["risks"][:5]], left=6.3, top=1.2, width=3.0, height=4.9)

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PII Exposure"
    slide.shapes.add_picture(visuals["pii"], Inches(0.6), Inches(1.2), width=Inches(5.6))
    add_bullets(slide, [f"{item['table']}.{item['column']}: {item['masking']}" for item in context["pii_findings"][:5]], left=6.3, top=1.2, width=3.0, height=4.9)

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Data Quality"
    slide.shapes.add_picture(visuals["dq"], Inches(0.6), Inches(1.2), width=Inches(5.6))
    add_bullets(slide, [f"{item.get('table')}: {item.get('issue')}" for item in context["data_quality"]["issues"][:5]], left=6.3, top=1.2, width=3.0, height=4.9)

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "AI Insights"
    add_bullets(slide, [item["summary"] for item in context["ai_insights"][:5]])

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Recommendations"
    add_bullets(slide, [item["action"] for item in context["recommendations"][:5]])

    presentation.save(path)
    return path


def export_pdf(job_id: str, context: Dict[str, Any], visuals: Dict[str, str]) -> str:
    path = os.path.join(REPORT_DIR, f"scan_report_{job_id}.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4, topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("BodySmall", parent=styles["BodyText"], fontSize=9, leading=12)
    story: List[Any] = []

    story.append(Paragraph("Data Governance Report", styles["Title"]))
    story.append(Spacer(1, 10))
    story.append(_section_heading("Executive Summary", styles))
    story.append(Paragraph(context["summary"]["executive_narrative"], styles["BodyText"]))
    story.append(Spacer(1, 8))

    def add_image_section(title: str, image_path: str) -> None:
        story.append(_section_heading(title, styles))
        story.append(Image(image_path, width=6.5 * inch, height=3.8 * inch))
        story.append(Spacer(1, 6))

    add_image_section("Risk Overview", visuals["risk"])
    risk_rows = [["Risk", "Assigned Owner", "Impact"]] + [
        [
            _paragraph_cell(item["summary"], body_style),
            _paragraph_cell(item["assigned_owner"], body_style),
            _paragraph_cell(item["impact"], body_style),
        ]
        for item in context["risks"][:8]
    ]
    risk_table = Table(risk_rows, colWidths=[2.3 * inch, 1.6 * inch, 2.4 * inch])
    risk_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eaf1f5")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cad5df")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
    ]))
    story.append(risk_table)
    story.append(Spacer(1, 8))

    add_image_section("PII Exposure", visuals["pii"])
    pii_rows = [["Table", "Column", "Exposure", "Recommendation"]] + [
        [
            _paragraph_cell(item["table"], body_style),
            _paragraph_cell(item["column"], body_style),
            _paragraph_cell(item["masking"], body_style),
            _paragraph_cell(item["recommendation"], body_style),
        ]
        for item in context["pii_findings"][:10]
    ]
    pii_table = Table(pii_rows, colWidths=[1.8 * inch, 1.2 * inch, 1.5 * inch, 2.0 * inch])
    pii_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eaf1f5")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cad5df")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
    ]))
    story.append(pii_table)
    story.append(Spacer(1, 8))

    add_image_section("Data Quality", visuals["dq"])
    story.append(Paragraph(context["data_quality"]["summary"], styles["BodyText"]))
    story.append(Spacer(1, 8))

    story.append(_section_heading("AI Insights", styles))
    for item in context["ai_insights"][:5] or [{"summary": "No AI insights were generated."}]:
        story.append(Paragraph(f"- {item.get('summary')}", body_style))
    story.append(Spacer(1, 8))

    story.append(_section_heading("Recommendations", styles))
    for group_name in ("critical", "improvements", "observations"):
        items = context["priorities"][group_name]
        if not items:
            continue
        story.append(Paragraph(group_name.title(), styles["Heading3"]))
        for item in items[:6]:
            story.append(Paragraph(f"- {item['summary']}<br/>Impact: {item['impact']}", body_style))
        story.append(Spacer(1, 4))

    for item in context["recommendations"][:6]:
        story.append(Paragraph(f"- {item['action']}", body_style))

    doc.build(story)
    return path


def generate_export(format_name: str, job_id: str, status_payload: Dict[str, Any]) -> Tuple[str, str]:
    context = build_report_context(status_payload, job_id)
    visuals = generate_visuals(context)

    if format_name == "json":
        path = export_raw_json(job_id, status_payload)
        return path, "application/json"
    if format_name == "pdf":
        path = export_pdf(job_id, context, visuals)
        return path, "application/pdf"
    if format_name == "pptx":
        path = export_pptx(job_id, context, visuals)
        return path, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if format_name == "docx":
        path = export_docx(job_id, context, visuals)
        return path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if format_name == "xlsx":
        path = export_xlsx(job_id, context, visuals)
        return path, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if format_name == "md":
        path = export_markdown(job_id, context)
        return path, "text/markdown"
    if format_name == "txt":
        path = export_text(job_id, context)
        return path, "text/plain"
    raise ValueError(f"Unsupported export format: {format_name}")
